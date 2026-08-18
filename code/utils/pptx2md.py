#!/usr/bin/env python3
"""Convert the UG lecture decks in slides/ into the site's per-week MDX decks.

Source of truth is `slides/**/{week}-UG-{Title}.pptx` (week 1 is spelled
`UG1-{Title}.pptx`). Lab decks (`*-UG-Lab*.pptx`) are out of scope — they feed
`labs_hw/` via the lab-author skill. Weeks past 13 are ignored.

The site keeps **one deck per week**, `site/content/decks/week-NN/slides.mdx`.
When several UG pptx map to the same week (`10-UG-Classifiers.pptx` and
`10-UG-Classifiers2.pptx`), the newest by mtime wins; the losers are reported so
the drop is visible rather than silent.

Why python-pptx and not raw XML
-------------------------------
python-pptx *is* a thin object layer over `ppt/slides/slideN.xml`, and it
resolves the two things that are tedious by hand: relationship parts (images,
layouts, notes) and layout inheritance. Auditing the corpus, the only content
raw XML exposes that python-pptx cannot reach is SmartArt (`ppt/diagrams/`) —
2 diagrams in 827 slides. So this module uses python-pptx for everything and
drops to `zipfile` + regex only for that SmartArt gap (`smartart_text`).

What earns the fidelity here is not the parser but *using signal the previous
version discarded*: slide-layout names, placeholder roles, shape geometry,
run-level bold/italic/monospace, hyperlinks, and speaker notes.

Hand-authored `key=`, `aiNotes=` and `<Quiz>` blocks in an existing deck are
harvested before regeneration and re-attached to slides whose titles still
match, so re-converting never destroys pedagogy.

Usage:
    python code/utils/pptx2md.py                  # all weeks -> build/decks-staging/
    python code/utils/pptx2md.py --promote        # ... and write into site/content/
    python code/utils/pptx2md.py FILE.pptx        # one deck (forces its week)
"""
import re
import sys
import html
import hashlib
import pathlib
import zipfile
import argparse
import collections

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE as ST

ROOT = pathlib.Path(__file__).resolve().parents[2]
SLIDES = ROOT / "slides"
SITE = ROOT / "site"
DECKS = SITE / "content" / "decks"
WEEKS = SITE / "content" / "weeks"
COURSE = SITE / "content" / "course.yaml"
PUBLIC = SITE / "public" / "decks"
STAGING = ROOT / "build" / "decks-staging"

MAX_WEEK = 13

# `2-UG-Business Data.pptx` -> (2, "Business Data");  `UG1-Intro.pptx` -> (1, "Intro")
UG_RE = re.compile(r"^(\d+)-UG-(.+)$", re.I)
UG1_RE = re.compile(r"^UG(\d+)-(.+)$", re.I)
LAB_RE = re.compile(r"^lab", re.I)

MONO_FONTS = {"consolas", "courier new", "menlo", "monaco", "cascadia mono", "lucida console"}

EMU_PER_IN = 914400
# An image displaying smaller than this is slide chrome (logo, badge), not content.
MIN_IMAGE_AREA_IN2 = 1.5
# An identical image appearing on this many slides is a template element.
CHROME_REPEAT = 3
# aiNotes feed a small on-device model — keep the budget tight.
AI_NOTES_MAX = 400


# --------------------------------------------------------------------------- #
# string helpers
# --------------------------------------------------------------------------- #
def slugify(s):
    s = re.sub(r"[^a-z0-9]+", "-", s.strip().lower())
    return re.sub(r"-+", "-", s).strip("-")


def collapse(s):
    return re.sub(r"\s+", " ", (s or "")).strip()


def js_str(s):
    """A double-quoted JS string literal, safe inside an MDX `{…}` expression.

    `<` and `{` need no escaping inside a string literal, which is why every
    text payload here rides in a prop rather than in MDX flow content."""
    return '"' + collapse(s).replace("\\", "\\\\").replace('"', '\\"') + '"'


def yaml_str(s):
    return '"' + collapse(s).replace("\\", "\\\\").replace('"', '\\"') + '"'


def attr(name, value):
    """`name={"…"}` — a JSX *expression container*, not a quoted attribute.

    A plain `name="…"` attribute has no escape mechanism in JSX, so any quote in
    the slide text breaks the parse. Wrapping in braces makes it a JS string,
    where `js_str`'s escaping is valid."""
    return f"{name}={{{js_str(value)}}}"


def unescape_js(s):
    """Inverse of `js_str`'s escaping, for values harvested out of an old deck."""
    return re.sub(r'\\(.)', r"\1", s or "")


def norm_title(s):
    """Matching key for carrying hand-authored props across a re-conversion."""
    return re.sub(r"[^a-z0-9]+", "", (s or "").lower())


# --------------------------------------------------------------------------- #
# run-level inline formatting  (**bold**, *italic*, `code`)
# --------------------------------------------------------------------------- #
def run_markup(run):
    """One text run as an Inline-compatible string. `Inline` supports exactly
    bold/italic/code — hyperlinks are hoisted to <Cta> by the caller."""
    text = run.text
    if not text.strip():
        return text
    font = run.font
    name = (font.name or "").lower()
    lead = len(text) - len(text.lstrip())
    trail = len(text) - len(text.rstrip())
    core = text.strip()
    if name in MONO_FONTS:
        core = f"`{core}`"
    else:
        if font.bold:
            core = f"**{core}**"
        if font.italic:
            core = f"*{core}*"
    return text[:lead] + core + text[len(text) - trail:] if trail else text[:lead] + core


def para_markup(para):
    """A paragraph's runs merged, coalescing adjacent identically-styled runs so
    PowerPoint's habit of splitting a word across runs doesn't yield `**a****b**`."""
    out = "".join(run_markup(r) for r in para.runs)
    for mark in ("**", "*", "`"):
        out = out.replace(mark + mark, "")
    return collapse(out)


def para_links(para):
    """[(text, url)] for runs carrying a hyperlink."""
    links = []
    for r in para.runs:
        addr = getattr(r.hyperlink, "address", None)
        if addr:
            links.append((collapse(r.text) or addr, addr))
    return links


def bullet_items(tf):
    """Text frame -> [(level, markup)], 1-based level, blanks dropped."""
    items = []
    for p in tf.paragraphs:
        text = para_markup(p)
        if text:
            items.append((min(p.level + 1, 4), text))
    return items


# --------------------------------------------------------------------------- #
# SmartArt — the one thing python-pptx cannot read, so go to the XML
# --------------------------------------------------------------------------- #
def smartart_text(pptx_path):
    """{diagram-part-name: [text, …]} pulled straight from `ppt/diagrams/dataN.xml`.

    python-pptx surfaces SmartArt as an opaque GraphicFrame; the node text only
    exists in the diagram data part, so this reaches into the zip for it."""
    out = {}
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if not re.match(r"ppt/diagrams/data\d+\.xml$", name):
                continue
            xml = z.read(name).decode("utf-8", "replace")
            texts = [collapse(html.unescape(t)) for t in re.findall(r"<a:t>(.*?)</a:t>", xml, re.S)]
            seen, uniq = set(), []
            for t in texts:
                if t and t not in seen:
                    seen.add(t)
                    uniq.append(t)
            if uniq:
                out[name] = uniq
    return out


# --------------------------------------------------------------------------- #
# geometry & shape classification
# --------------------------------------------------------------------------- #
def area_in2(sh):
    return ((sh.width or 0) / EMU_PER_IN) * ((sh.height or 0) / EMU_PER_IN)


def reading_order(shapes):
    return sorted(shapes, key=lambda s: ((s.top or 0), (s.left or 0)))


def is_title(sh):
    if not getattr(sh, "is_placeholder", False):
        return False
    try:
        return "TITLE" in str(sh.placeholder_format.type)
    except (ValueError, AttributeError):
        return False


def image_fingerprints(prs):
    """sha1 -> number of slides the image appears on, for chrome detection."""
    counts = collections.Counter()
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.shape_type == ST.PICTURE:
                try:
                    counts[hashlib.sha1(sh.image.blob).hexdigest()] += 1
                except Exception:
                    pass
    return counts


# --------------------------------------------------------------------------- #
# MDX block emission
# --------------------------------------------------------------------------- #
def bullets_block(items, indent=""):
    out = [indent + "<Bullets items={["]
    for level, text in items:
        out.append(f'{indent}  {{ l: {level}, h: {js_str(text)} }},')
    out.append(indent + "]} />")
    return out


def datatable_block(columns, rows, indent=""):
    cols = ", ".join(js_str(c) for c in columns)
    out = [indent + "<DataTable", f"{indent}  columns={{[{cols}]}}", indent + "  rows={["]
    for r in rows:
        out.append(f'{indent}    [' + ", ".join(js_str(c) for c in r) + "],")
    out += [indent + "  ]}", indent + "/>"]
    return out


def image_block(src, alt, indent=""):
    return [f'{indent}<ImageCaption {attr("src", src)} {attr("alt", alt)} />']


def cta_block(label, href, indent=""):
    return [f'{indent}<Cta {attr("label", label)} {attr("href", href)} />']


def table_data(tbl):
    rows = [[collapse(c.text) for c in r.cells] for r in tbl.rows]
    if not rows:
        return None
    keep = [i for i in range(len(rows[0])) if any(r[i] for r in rows)]
    rows = [[r[i] for i in keep] for r in rows]
    return (rows[0], rows[1:]) if rows and rows[0] else None


# --------------------------------------------------------------------------- #
# slide -> structured blocks
# --------------------------------------------------------------------------- #
class Block:
    """One emitted unit, tagged with the geometry used for column inference.

    Text blocks keep their `items` so stacked text boxes can be merged before
    rendering — PowerPoint often splits one logical list across several boxes."""

    def __init__(self, lines, left, right, kind, items=None):
        self.lines = lines
        self.left = left
        self.right = right
        self.kind = kind
        self.items = items


def merge_stacked_text(blocks):
    """Fold consecutive text blocks sharing a column into a single <Bullets>."""
    out = []
    for b in blocks:
        prev = out[-1] if out else None
        if prev and prev.kind == b.kind == "text" and _overlaps(prev, b):
            prev.items += b.items
            prev.left, prev.right = min(prev.left, b.left), max(prev.right, b.right)
            prev.lines = bullets_block(prev.items)
        else:
            out.append(b)
    return out


def _overlaps(a, b, frac=0.6):
    """True if two blocks share most of their horizontal extent (same column)."""
    span = min(a.right, b.right) - max(a.left, b.left)
    return span > 0 and span >= frac * min(a.right - a.left, b.right - b.left)


def collect_blocks(slide, ctx):
    """Content blocks for one slide, in reading order, plus any hyperlinks."""
    blocks, links = [], []
    body = reading_order(s for s in slide.shapes if not is_title(s))
    for sh in body:
        left = sh.left or 0
        right = left + (sh.width or 0)
        try:
            if sh.shape_type == ST.PICTURE:
                digest = hashlib.sha1(sh.image.blob).hexdigest()
                if ctx["chrome"][digest] >= CHROME_REPEAT or area_in2(sh) < MIN_IMAGE_AREA_IN2:
                    continue  # template logo / decorative badge
                name = ctx["save_image"](sh, digest)
                alt = sh.name if sh.name and not sh.name.lower().startswith(("picture", "image")) else ctx["slide_title"]
                blocks.append(Block(image_block(f'{ctx["img_base"]}/{name}', alt or "slide image"), left, right, "image"))
            elif sh.has_table:
                data = table_data(sh.table)
                if data:
                    blocks.append(Block(datatable_block(*data), left, right, "table"))
            elif sh.shape_type == ST.PLACEHOLDER and getattr(sh, "has_chart", False):
                continue
            elif sh.has_text_frame and sh.text_frame.text.strip():
                items = bullet_items(sh.text_frame)
                if items:
                    blocks.append(Block(bullets_block(items), left, right, "text", items))
                for p in sh.text_frame.paragraphs:
                    links += para_links(p)
        except Exception as e:
            print(f"    ! shape skipped ({sh.shape_type}): {e}", file=sys.stderr)
    return blocks, links


def two_column(blocks, slide_width):
    """Re-emit exactly two side-by-side blocks as <TwoColumn>, else None.

    PowerPoint encodes 'these two things are parallel' purely as geometry; this
    is the one structural cue worth recovering automatically."""
    if len(blocks) != 2 or not slide_width:
        return None
    a, b = sorted(blocks, key=lambda x: x.left)
    if a.right > slide_width * 0.58 or b.left < slide_width * 0.42:
        return None
    ratio = "1-1"
    span_a, span_b = a.right - a.left, b.right - b.left
    if span_b and span_a / span_b > 1.6:
        ratio = "2-1"
    elif span_a and span_b / span_a > 1.6:
        ratio = "1-2"
    out = [f'<TwoColumn ratio="{ratio}">', '  <div slot="left">']
    out += ["  " + ln for ln in a.lines]
    out += ["  </div>", '  <div slot="right">']
    out += ["  " + ln for ln in b.lines]
    out += ["  </div>", "</TwoColumn>"]
    return out


def notes_text(slide):
    if not slide.has_notes_slide:
        return ""
    text = collapse(slide.notes_slide.notes_text_frame.text)
    if len(text) < 15:
        return ""
    if len(text) > AI_NOTES_MAX:
        cut = text[:AI_NOTES_MAX].rsplit(". ", 1)[0]
        text = (cut + ".") if len(cut) > AI_NOTES_MAX * 0.5 else text[:AI_NOTES_MAX].rstrip() + "…"
    return text


# --------------------------------------------------------------------------- #
# harvest hand-authored props from an existing deck
# --------------------------------------------------------------------------- #
def scan_tag(text, start):
    """End index of the tag opening at `start`, ignoring `>` inside quotes/braces."""
    i, depth, quote = start, 0, None
    while i < len(text):
        c = text[i]
        if quote:
            if c == quote:
                quote = None
        elif c in "\"'":
            quote = c
        elif c == "{":
            depth += 1
        elif c == "}":
            depth -= 1
        elif c == ">" and depth == 0:
            return i + 1
        i += 1
    return len(text)


def attr_value(head, name):
    """Read `name="…"` or `name={"…"}` off a tag. Both spellings exist in the
    decks: hand-authored slides use the plain attribute, this script emits the
    expression container (see `attr`)."""
    m = re.search(r'\b%s=(?:"((?:[^"\\]|\\.)*)"|\{\s*"((?:[^"\\]|\\.)*)"\s*\})' % name, head)
    if not m:
        return None
    return m.group(1) if m.group(1) is not None else m.group(2)


def harvest(path):
    """{normalized slide title: [{'key','aiNotes','quizzes'}, …]} from a deck."""
    if not path.exists():
        return {}
    text = path.read_text()
    found = collections.defaultdict(list)
    for m in re.finditer(r"<Slide\b", text):
        head_end = scan_tag(text, m.start())
        head = text[m.start():head_end]
        title = attr_value(head, "title")
        if title is None:
            continue
        close = text.find("</Slide>", head_end)
        body = text[head_end:close if close != -1 else len(text)]
        quizzes = []
        for q in re.finditer(r"<Quiz\b", body):
            quizzes.append(body[q.start():scan_tag(body, q.start())].strip())
        km, am = attr_value(head, "key"), attr_value(head, "aiNotes")
        if km or am or quizzes:
            found[norm_title(title)].append(
                {"title": title, "key": km, "aiNotes": am, "quizzes": quizzes}
            )
    return found


def orphan_report(keep):
    """Whatever `harvest` found but `build_deck` never re-attached — a title
    changed upstream, so the authored material needs a manual home."""
    out = []
    for entries in keep.values():
        for e in entries:
            lines = [f'## {e["title"]}']
            if e.get("key"):
                lines.append(f'key="{e["key"]}"')
            if e.get("aiNotes"):
                lines.append(f'aiNotes="{e["aiNotes"]}"')
            lines += e.get("quizzes", [])
            out.append("\n".join(lines))
    return out


# --------------------------------------------------------------------------- #
# deck assembly
# --------------------------------------------------------------------------- #
def build_deck(pptx_path, week_id, week_num, out_root, keep):
    prs = Presentation(pptx_path)
    slug = slugify(pptx_path.stem)
    img_dir = out_root / "public" / "decks" / week_id / slug
    saved = {}

    def save_image(sh, digest):
        if digest in saved:
            return saved[digest]
        img = sh.image
        img_dir.mkdir(parents=True, exist_ok=True)
        name = f"img{len(saved) + 1}.{img.ext}"
        (img_dir / name).write_bytes(img.blob)
        saved[digest] = name
        return name

    diagrams = list(smartart_text(pptx_path).values())
    ctx = {
        "chrome": image_fingerprints(prs),
        "save_image": save_image,
        "img_base": f"/decks/{week_id}/{slug}",
        "smartart": diagrams,
        "slide_title": "",
    }

    slides = list(prs.slides)
    deck_title = collapse(UG_RE.sub(r"\2", pptx_path.stem)) or pptx_path.stem
    subtitle = ""
    if slides:
        t = slides[0].shapes.title
        if t and t.text.strip():
            deck_title = collapse(t.text)
        for sh in reading_order(s for s in slides[0].shapes if not is_title(s)):
            if sh.has_text_frame and sh.text_frame.text.strip():
                subtitle = collapse(sh.text_frame.text)
                break

    used, body, reattached = set(), [], collections.Counter()

    # ---- cover ----
    icon = week_icon(week_id)
    cover = f'<Slide variant="title" kicker="Week {week_num}"'
    if icon:
        cover += f' icon="{icon}"'
    body += [cover + ">", ""]
    body.append(deck_title if not subtitle else subtitle)
    body += ["", f'<CoverMeta instructors={{["Maclean Gaulin"]}} />', "", "</Slide>", ""]
    used |= {"Slide", "CoverMeta"}

    diagram_queue = list(diagrams)
    for slide in slides[1:]:
        layout = (slide.slide_layout.name or "").strip()
        t = slide.shapes.title
        title = collapse(t.text) if t and t.text.strip() else ""
        ctx["slide_title"] = title or deck_title

        blocks, links = collect_blocks(slide, ctx)

        # SmartArt lives in a GraphicFrame python-pptx can't read; splice in the
        # node text recovered from ppt/diagrams/ instead.
        if diagram_queue and _has_diagram(slide):
            nodes = [(1, n) for n in diagram_queue.pop(0)]
            blocks.append(Block(bullets_block(nodes), 0, 0, "text", nodes))

        blocks = merge_stacked_text(blocks)
        if not title and not blocks:
            continue

        # A mid-deck "Title Slide" layout is a section divider.
        if layout.lower() == "title slide" and title:
            body += [f'<Slide variant="title" {attr("title", title)}>', "", "</Slide>", ""]
            continue

        lines = two_column(blocks, prs.slide_width)
        if lines:
            used.add("TwoColumn")
        else:
            lines = []
            for b in blocks:
                lines += b.lines + [""]
            lines = lines[:-1] if lines else lines
        for b in blocks:
            used.add({"text": "Bullets", "table": "DataTable", "image": "ImageCaption"}[b.kind])

        attrs = []
        if title:
            attrs.append(attr("title", title))

        carried = keep.get(norm_title(title))
        entry = carried.pop(0) if carried else None
        if entry and entry.get("key"):
            attrs.append(attr("key", unescape_js(entry["key"])))
            reattached["key carried"] += 1
        if entry and entry.get("aiNotes"):
            attrs.append(attr("aiNotes", unescape_js(entry["aiNotes"])))
            reattached["aiNotes carried"] += 1
        elif notes_text(slide):
            attrs.append(attr("aiNotes", notes_text(slide)))
            reattached["aiNotes from speaker notes"] += 1

        body.append(f"<Slide {' '.join(attrs)}>" if attrs else "<Slide>")
        if lines:
            body += [""] + lines
        for label, url in dict(links).items():
            body += [""] + cta_block(label, url)
            used.add("Cta")
        for quiz in (entry or {}).get("quizzes", []):
            body += [""] + quiz.splitlines()
            used.add("Quiz")
            reattached["quiz carried"] += 1
        body += ["", "</Slide>", ""]

    order = ["Slide", "CoverMeta", "Bullets", "TwoColumn", "DataTable", "ImageCaption", "Cta", "Quiz"]
    imports = ", ".join(c for c in order if c in used)
    fm = ["---", f"week: {week_id}", f"title: {yaml_str(deck_title)}"]
    if subtitle and subtitle != deck_title:
        fm.append(f"subtitle: {yaml_str(subtitle)}")
    fm += ["---", "", f"import {{ {imports} }} from '@components/blocks';", ""]
    return "\n".join(fm + body).rstrip() + "\n", reattached


DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def _has_diagram(slide):
    """True if any shape is a SmartArt graphic frame (identified by graphicData uri)."""
    from lxml import etree

    return any(
        el.get("uri") == DIAGRAM_URI
        for el in slide.shapes._spTree.iter()
        if isinstance(el.tag, str) and etree.QName(el).localname == "graphicData"
    )


def week_icon(week_id):
    path = WEEKS / f"{week_id}.yaml"
    if not path.exists():
        return None
    m = re.search(r'^icon:\s*"?([A-Za-z]+)"?\s*$', path.read_text(), re.M)
    return m.group(1) if m else None


# --------------------------------------------------------------------------- #
# selection
# --------------------------------------------------------------------------- #
def parse_ug(path):
    """(week, title) for an in-scope UG lecture deck, else None."""
    for rx in (UG_RE, UG1_RE):
        m = rx.match(path.stem)
        if m:
            week, title = int(m.group(1)), m.group(2)
            if LAB_RE.match(title) or week > MAX_WEEK:
                return None
            return week, title
    return None


def select_decks():
    """One winning pptx per week (newest mtime), plus the candidates it beat."""
    by_week = collections.defaultdict(list)
    for p in sorted(SLIDES.rglob("*.pptx")):
        hit = parse_ug(p)
        if hit:
            by_week[hit[0]].append(p)
    chosen, dropped = {}, {}
    for week, paths in by_week.items():
        paths.sort(key=lambda p: p.stat().st_mtime, reverse=True)
        chosen[week] = paths[0]
        dropped[week] = paths[1:]
    return chosen, dropped


def ensure_week_stub(week_id, num, title, out_root):
    path = out_root / "weeks" / f"{week_id}.yaml"
    if path.exists() or (WEEKS / f"{week_id}.yaml").exists():
        return
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(f"number: {num}\ntitle: {yaml_str(title)}\nsummary: {yaml_str(title)}\n")
    print(f"  + week stub {path.name}")


# --------------------------------------------------------------------------- #
# driver
# --------------------------------------------------------------------------- #
def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("files", nargs="*", help="specific pptx to convert (default: all UG lecture decks)")
    ap.add_argument("--promote", action="store_true", help="write into site/content/ instead of the staging tree")
    args = ap.parse_args(argv)

    out_root = SITE / "content" if args.promote else STAGING
    pub_root = SITE if args.promote else STAGING

    if args.files:
        targets, dropped = {}, {}
        for f in args.files:
            p = pathlib.Path(f).resolve()
            hit = parse_ug(p)
            if not hit:
                print(f"! not an in-scope UG lecture deck: {p.name}", file=sys.stderr)
                continue
            targets[hit[0]] = p
    else:
        targets, dropped = select_decks()

    for week in sorted(dropped or {}):
        for p in dropped[week]:
            print(f"  ~ week {week}: skipped {p.name} (older than {targets[week].name})")

    total = collections.Counter()
    for week in sorted(targets):
        pptx = targets[week]
        week_id = f"week-{week:02d}"
        deck_path = out_root / "decks" / week_id / "slides.mdx"
        keep = harvest(DECKS / week_id / "slides.mdx")
        print(f"{week_id}  <-  {pptx.name}")

        mdx, reattached = build_deck(pptx, week_id, week, pub_root, keep)
        deck_path.parent.mkdir(parents=True, exist_ok=True)
        deck_path.write_text(mdx)

        note = "; ".join(f"{v} {k}" for k, v in sorted(reattached.items()) if v)
        print(f"  -> {deck_path.relative_to(out_root.parent)}" + (f"   [{note}]" if note else ""))

        orphans = orphan_report(keep)
        if orphans:
            path = STAGING / f"{week_id}.orphans.md"
            path.parent.mkdir(parents=True, exist_ok=True)
            header = (f"# {week_id} — authored content with no matching slide\n\n"
                      f"These came from the previous slides.mdx but no regenerated slide title\n"
                      f"matched. Re-home them by hand, then delete this file.\n\n")
            path.write_text(header + "\n\n".join(orphans) + "\n")
            print(f"  ! {len(orphans)} orphaned slide(s) -> {path.name}")
            total["orphaned"] += len(orphans)

        ensure_week_stub(week_id, week, collapse(parse_ug(pptx)[1]), out_root)
        total.update(reattached)

    where = "site/content" if args.promote else STAGING.relative_to(ROOT)
    print(f"\nDone. {len(targets)} week(s) -> {where}.")
    if total:
        print("Carried forward: " + ", ".join(f"{v} {k}" for k, v in sorted(total.items())))
    if not args.promote:
        print("Review, then re-run with --promote to write into site/content/.")


if __name__ == "__main__":
    main()
